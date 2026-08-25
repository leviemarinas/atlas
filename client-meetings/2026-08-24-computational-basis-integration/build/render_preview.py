# -*- coding: utf-8 -*-
"""Render the generated PPTX to PNGs by walking its shapes.

There is no PowerPoint or LibreOffice on this machine, so this reads the file
back with python-pptx and draws every shape with PIL. Font metrics differ
slightly from PowerPoint, but it is faithful enough to catch the things that
actually go wrong: overlapping boxes, text that overruns its frame, and images
placed outside their intended area.
"""
import os
import sys

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

SCALE = 110  # px per inch
EMU_IN = 914400

FONTS = {
    ("Segoe UI", False): r"C:\Windows\Fonts\segoeui.ttf",
    ("Segoe UI", True): r"C:\Windows\Fonts\segoeuib.ttf",
    ("Segoe UI Semibold", False): r"C:\Windows\Fonts\seguisb.ttf",
    ("Segoe UI Semibold", True): r"C:\Windows\Fonts\seguisb.ttf",
    ("Consolas", False): r"C:\Windows\Fonts\consola.ttf",
    ("Consolas", True): r"C:\Windows\Fonts\consolab.ttf",
}
_cache = {}


def font(name, size, bold=False, italic=False):
    key = (name, bold, round(size))
    if key not in _cache:
        path = FONTS.get((name, bold)) or FONTS[("Segoe UI", bold)]
        _cache[key] = ImageFont.truetype(path, max(6, round(size)))
    return _cache[key]


def px(emu):
    return round(emu / EMU_IN * SCALE)


def rgb(color, default=(0, 0, 0)):
    try:
        if color and color.type is not None and color.rgb is not None:
            return tuple(color.rgb)
    except Exception:
        pass
    return default


def wrap(draw, text, fnt, width):
    lines, line = [], ""
    for word in text.split():
        probe = (line + " " + word).strip()
        if draw.textlength(probe, font=fnt) <= width or not line:
            line = probe
        else:
            lines.append(line)
            line = word
    if line:
        lines.append(line)
    return lines


def draw_text_frame(draw, shape, overflow):
    tf = shape.text_frame
    x0, y0 = px(shape.left), px(shape.top)
    w, h = px(shape.width), px(shape.height)
    ml = px(tf.margin_left or 0)
    mr = px(tf.margin_right or 0)
    mt = px(tf.margin_top or 0)
    inner = max(10, w - ml - mr)
    y = y0 + mt
    total = 0
    blocks = []
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            total += 6
            blocks.append((None, None, None, 6, None))
            continue
        r = runs[0]
        size = (r.font.size.pt if r.font.size else 12)
        fnt = font(r.font.name or "Segoe UI", size * SCALE / 72,
                   bool(r.font.bold), bool(r.font.italic))
        text = "".join(run.text for run in runs)
        ls = p.line_spacing or 1.2
        lh = size * SCALE / 72 * (ls if ls > 0.2 else 1.2)
        before = (p.space_before.pt if p.space_before else 0) * SCALE / 72
        after = (p.space_after.pt if p.space_after else 0) * SCALE / 72
        lines = wrap(draw, text, fnt, inner)
        total += before + len(lines) * lh + after
        blocks.append((lines, fnt, rgb(r.font.color, (40, 40, 40)), lh, (before, after, p)))

    # vertical anchoring
    va = str(tf.vertical_anchor)
    if "MIDDLE" in va:
        y = y0 + (h - total) / 2
    elif "BOTTOM" in va:
        y = y0 + h - total - mt

    for block in blocks:
        lines, fnt, color, lh, meta = block
        if lines is None:
            y += lh
            continue
        before, after, p = meta
        y += before
        for ln in lines:
            tw = draw.textlength(ln, font=fnt)
            align = str(p.alignment)
            lx = x0 + ml
            if "CENTER" in align:
                lx = x0 + (w - tw) / 2
            elif "RIGHT" in align:
                lx = x0 + w - mr - tw
            draw.text((lx, y), ln, font=fnt, fill=color)
            y += lh
        y += after

    if total > h + 4 and shape.has_text_frame and any(b[0] for b in blocks):
        overflow.append((shape.shape_id, round(total - h), (x0, y0, w, h),
                         " ".join(l for b in blocks if b[0] for l in b[0])[:70]))


def render(path, outdir, limit=None):
    prs = Presentation(path)
    W, H = px(prs.slide_width), px(prs.slide_height)
    os.makedirs(outdir, exist_ok=True)
    problems = []
    for i, slide in enumerate(prs.slides, 1):
        if limit and i not in limit:
            continue
        im = Image.new("RGB", (W, H), (255, 255, 255))
        d = ImageDraw.Draw(im)
        overflow = []
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pic = Image.open(shape.image.blob and __import__("io").BytesIO(shape.image.blob))
                    pic = pic.convert("RGB").resize((max(1, px(shape.width)), max(1, px(shape.height))))
                    im.paste(pic, (px(shape.left), px(shape.top)))
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    fill = None
                    try:
                        if shape.fill.type is not None and str(shape.fill.type) != "MSO_FILL_TYPE.BACKGROUND (5)":
                            fill = rgb(shape.fill.fore_color, None)
                    except Exception:
                        fill = None
                    line = None
                    try:
                        line = tuple(shape.line.color.rgb)
                    except Exception:
                        line = None
                    box = (px(shape.left), px(shape.top),
                           px(shape.left) + px(shape.width), px(shape.top) + px(shape.height))
                    if "OVAL" in str(shape.shape_type) or shape.auto_shape_type is not None and "OVAL" in str(shape.auto_shape_type):
                        if fill:
                            d.ellipse(box, fill=fill, outline=line)
                    else:
                        d.rounded_rectangle(box, radius=6, fill=fill, outline=line)
                if shape.has_text_frame and shape.text_frame.text.strip():
                    draw_text_frame(d, shape, overflow)
            except Exception as exc:
                problems.append((i, f"shape error: {exc}"))
        im.save(os.path.join(outdir, f"slide-{i:02d}.png"))
        for o in overflow:
            problems.append((i, f"overflow +{o[1]}px  '{o[3]}'"))
    return problems


if __name__ == "__main__":
    src = sys.argv[1] if len(sys.argv) > 1 else os.path.join(
        "..", "outputs", "ATLAS_2026-08-24_Computational_Basis_Integration_PreMeeting_v03.pptx")
    probs = render(src, "preview")
    for slide, msg in probs:
        print(f"slide {slide:>2}: {msg}")
    print(f"\n{len(probs)} issues")
