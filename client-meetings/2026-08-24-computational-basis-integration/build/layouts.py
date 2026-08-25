# -*- coding: utf-8 -*-
"""Reusable slide layouts.

Every layout takes the same shape — eyebrow, heading, subhead, body, footnote —
so the deck reads as one document rather than a pile of one-off pages.
"""
import os

from PIL import Image
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from theme import (BODY_FONT, CANVAS, CONTENT_W, DECIDE, DECIDE_WASH, H, HEAD_FONT, INK,
                   fit_size, text_lines,
                   LINE, LIVE, LIVE_WASH, MARGIN, MARK, MARK_WASH, MONO_FONT, MUTED, VIOLET,
                   VIOLET_2, W, WASH, WHITE, badge, blank, circle, eyebrow, footnote,
                   heading, page_bg, para, path_chip, picture_fit, rect, subhead,
                   textbox)

ASSETS = "deck-assets"


def asset(name):
    return os.path.join(ASSETS, name + ".png")


# ------------------------------------------------------------------- covers
def cover(prs, kicker, title, sub, chips, footer):
    s = blank(prs)
    page_bg(s, INK)
    rect(s, 0, 0, Inches(0.16), H, fill=VIOLET_2, radius=0)
    band = rect(s, W - Inches(4.3), 0, Inches(4.3), H, fill=None, radius=0)
    band.fill.solid()
    band.fill.fore_color.rgb = INK

    _, tf = textbox(s, MARGIN + Inches(0.3), Inches(1.55), Inches(8.4), Inches(0.3))
    para(tf, kicker.upper(), size=11.5, color=VIOLET_2, font=HEAD_FONT, bold=True, first=True,
         space_after=0)
    _, tf = textbox(s, MARGIN + Inches(0.3), Inches(2.02), Inches(8.6), Inches(1.5))
    para(tf, title, size=52, color=WHITE, font=HEAD_FONT, bold=True, first=True,
         space_after=0, line=1.02)
    _, tf = textbox(s, MARGIN + Inches(0.3), Inches(3.42), Inches(7.9), Inches(1.0))
    para(tf, sub, size=15, color=RGB_SOFT(), first=True, space_after=0, line=1.4)

    x = MARGIN + Inches(0.3)
    for i, chip in enumerate(chips):
        w = Inches(0.108 * len(chip) + 0.66)
        c = rect(s, x, Inches(4.62), w, Inches(0.46), fill=None, outline=VIOLET_2, radius=0.4)
        tfc = c.text_frame
        tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tfc.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{i + 1}   {chip}"
        r.font.size = Pt(11.5)
        r.font.name = HEAD_FONT
        r.font.bold = True
        r.font.color.rgb = WHITE
        x += w + Inches(0.16)

    _, tf = textbox(s, MARGIN + Inches(0.3), Inches(5.66), Inches(8.6), Inches(0.9))
    para(tf, footer, size=11, color=VIOLET_2, first=True, space_after=0, line=1.5)
    return s


def RGB_SOFT():
    from pptx.dml.color import RGBColor
    return RGBColor(0xC5, 0xB8, 0xDC)


def section(prs, number, title, sub, items):
    s = blank(prs)
    page_bg(s, INK)
    rect(s, 0, 0, Inches(0.16), H, fill=VIOLET_2, radius=0)
    _, tf = textbox(s, MARGIN + Inches(0.3), Inches(1.9), Inches(1.6), Inches(1.0))
    para(tf, f"{number:02d}", size=64, color=VIOLET_2, font=HEAD_FONT, bold=True, first=True,
         space_after=0, line=1.0)
    _, tf = textbox(s, MARGIN + Inches(0.3), Inches(2.86), Inches(9.6), Inches(1.0))
    para(tf, title, size=38, color=WHITE, font=HEAD_FONT, bold=True, first=True,
         space_after=0, line=1.06)
    _, tf = textbox(s, MARGIN + Inches(0.3), Inches(3.98), Inches(8.4), Inches(0.8))
    para(tf, sub, size=14.5, color=RGB_SOFT(), first=True, space_after=0, line=1.4)

    x = MARGIN + Inches(0.3)
    for item in items:
        _, tf = textbox(s, x, Inches(5.28), Inches(2.62), Inches(1.2))
        para(tf, item, size=11.5, color=WHITE, font=HEAD_FONT, bold=True, first=True,
             space_after=0, line=1.3)
        rect(s, x, Inches(5.08), Inches(2.34), Inches(0.035), fill=VIOLET_2, radius=0)
        x += Inches(2.9)
    return s


# --------------------------------------------------------------- page frame
def page(prs, eb, title, sub=None, note=None, path=None):
    """Standard page frame. When a navigation path is given it sits directly
    under the heading, and the subhead moves below it."""
    s = blank(prs)
    page_bg(s, WHITE)
    rect(s, 0, 0, W, Inches(0.10), fill=VIOLET, radius=0)
    eyebrow(s, eb)
    heading(s, title)
    if path:
        path_chip(s, path)
    if sub:
        subhead(s, sub, y=Inches(1.92) if path else Inches(1.42))
    if note:
        footnote(s, note)
    return s


def body_top(spec):
    """Where slide content starts, given whether the page has a path and a subhead."""
    has_path = bool(spec.get("path"))
    has_sub = bool(spec.get("sub"))
    if has_path and has_sub:
        return Inches(2.52)
    if has_path:
        return Inches(2.06)
    if has_sub:
        return Inches(2.02)
    return Inches(1.66)


def note_card(slide, x, y, w, h, num, title, body, num_color=MARK, num_fill=MARK,
              size_title=12.5, size_body=10.8):
    """A numbered note. Both texts are sized down until they fit, so a long
    sentence shrinks rather than spilling over the card edge."""
    card = rect(slide, x, y, w, h, fill=CANVAS, outline=LINE, radius=0.06)
    pad = Inches(0.22)
    d = Inches(0.30)
    circle(slide, x + pad + Inches(0.15), y + Inches(0.30), d, fill=num_fill, label=num, size=11)
    tx = x + pad + Inches(0.41)
    tw = w - (tx - x) - pad
    t_pt = fit_size(title, tw / 914400, 0.42, max_pt=size_title, min_pt=9.5,
                    line=1.15, font=HEAD_FONT, bold=True)
    t_lines = 2 if t_pt <= size_title - 1.5 else 1
    _, tf = textbox(slide, tx, y + Inches(0.16), tw, Inches(0.44))
    para(tf, title, size=t_pt, color=INK, font=HEAD_FONT, bold=True, first=True,
         space_after=0, line=1.15)

    by = y + Inches(0.14 + 0.30 * t_lines)
    bw = w - pad * 2
    bh = y + h - by - Inches(0.11)
    b_pt = fit_size(body, bw / 914400, bh / 914400, max_pt=size_body, min_pt=7.4, line=1.26)
    _, tf = textbox(slide, x + pad, by, bw, bh)
    para(tf, body, size=b_pt, color=MUTED, first=True, space_after=0, line=1.26)
    return card


def guide(prs, spec, num_fill=MARK):
    """Field-guide page: the annotated screen plus one note per marker.

    Wide screens run across the top with the notes in a row beneath; everything
    else puts the screen on the left and the notes down the right. The notes
    column starts where the image actually ends, so a tall narrow capture does
    not leave a band of dead space.
    """
    s = page(prs, spec["eyebrow"], spec["title"], spec.get("sub"), spec.get("note"),
             spec.get("path"))
    img = asset(spec["_key"])
    iw, ih = Image.open(img).size
    aspect = iw / ih
    marks = spec["marks"]
    labels = spec.get("labels") or [str(i + 1) for i in range(len(marks))]
    n = len(marks)
    top = body_top(spec)
    avail_h = H - top - Inches(0.72)
    gap = Inches(0.15)

    if aspect >= 2.6:
        # wide screen across the top, notes in one row beneath. The picture is
        # sized first so a readable screenshot is never traded for card padding.
        cols = n
        col_in = ((CONTENT_W - gap * (cols - 1)) // cols - Inches(0.44)) / 914400
        body_pt = 10.0 if cols <= 5 else 9.2
        need = max(text_lines(b, col_in, body_pt) for _, _, b in marks) * body_pt * 1.26 / 72
        band = Inches(min(2.35, max(1.30, need + 0.66)))
        pic_h = min(int(CONTENT_W / aspect), avail_h - band - Inches(0.28))
        slack = max(0, avail_h - pic_h - band - Inches(0.28))
        py = top + slack // 2
        picture_fit(s, img, MARGIN, py, CONTENT_W, pic_h, valign="top")
        cw = (CONTENT_W - gap * (cols - 1)) // cols
        cy = py + pic_h + Inches(0.28)
        for i, (_, t, b) in enumerate(marks):
            note_card(s, MARGIN + i * (cw + gap), cy, cw, band,
                      labels[i], t, b, num_fill=num_fill,
                      size_title=11.4 if cols <= 5 else 10.4,
                      size_body=body_pt)
    else:
        # five or more notes need two columns to stay readable, so the image is
        # capped to leave room for them rather than squeezing the text
        cap = Inches(7.3) if n < 5 else (W - MARGIN * 2 - Inches(0.32) - Inches(6.15))
        pic_w = min(cap, int(avail_h * aspect))
        picture_fit(s, img, MARGIN, top, pic_w, avail_h, align="left")
        cx = MARGIN + pic_w + Inches(0.32)
        col_w = W - cx - MARGIN
        cols = 2 if (col_w > Inches(6.0) and n >= 4) else 1
        rows = -(-n // cols)
        cw = (col_w - gap * (cols - 1)) // cols
        ch = (avail_h - gap * (rows - 1)) // rows
        for i, (_, t, b) in enumerate(marks):
            r, c = divmod(i, cols)
            note_card(s, cx + c * (cw + gap), top + r * (ch + gap), cw, ch,
                      labels[i], t, b, num_fill=num_fill,
                      size_title=11.6 if rows <= 4 else 10.8,
                      size_body=10.3 if rows <= 4 else 9.5)
    return s


def screen(prs, key, spec, caption_items=None):
    """A full-width screen with an optional row of short captions beneath.

    The screen and its captions are treated as one block and centred together,
    so a wide, short capture does not leave a band of empty slide under it.
    """
    s = page(prs, spec["eyebrow"], spec["title"], spec.get("sub"), spec.get("note"),
             spec.get("path"))
    top = body_top(spec)
    avail = H - top - Inches(0.74)
    cap_h = Inches(0.94) if caption_items else 0
    gap = Inches(0.34) if caption_items else 0

    iw, ih = Image.open(asset(key)).size
    box_h = avail - cap_h - gap
    pic_h = min(box_h, int(CONTENT_W * ih / iw))
    block = pic_h + gap + cap_h
    py = top + max(0, (avail - block) // 2)
    picture_fit(s, asset(key), MARGIN, py, CONTENT_W, pic_h, valign="top")

    if caption_items:
        n = len(caption_items)
        cgap = Inches(0.16)
        cw = (CONTENT_W - cgap * (n - 1)) // n
        cy = py + pic_h + gap
        for i, (t, b) in enumerate(caption_items):
            x = MARGIN + i * (cw + cgap)
            rect(s, x, cy, Inches(0.036), Inches(0.86), fill=VIOLET_2, radius=0)
            _, tf = textbox(s, x + Inches(0.16), cy, cw - Inches(0.2), Inches(0.3))
            para(tf, t, size=11.5, color=INK, font=HEAD_FONT, bold=True, first=True,
                 space_after=3, line=1.15)
            _, tf = textbox(s, x + Inches(0.16), cy + Inches(0.26), cw - Inches(0.2), Inches(0.68))
            para(tf, b, size=fit_size(b, (cw - Inches(0.2)) / 914400, 0.66,
                                      max_pt=10.0, min_pt=8.2, line=1.26),
                 color=MUTED, first=True, space_after=0, line=1.26)
    return s


def cards(prs, eb, title, sub, items, cols=3, note=None, accent=VIOLET,
          wash=WASH, numbered=True, body_size=10.8, path=None):
    """A grid of titled cards. Row height follows the longest body in the grid,
    and the block is centred, so short copy does not leave tall empty cards."""
    s = page(prs, eb, title, sub, note, path)
    top = body_top({"path": path, "sub": sub}) + Inches(0.04)
    rows = (len(items) + cols - 1) // cols
    gap = Inches(0.18)
    cw = (CONTENT_W - gap * (cols - 1)) // cols
    avail = H - top - Inches(0.74)
    max_ch = (avail - gap * (rows - 1)) // rows

    body_w = (cw - Inches(0.52)) / 914400
    need = max(text_lines(b, body_w, body_size) for _, b in items) * body_size * 1.3 / 72
    ch = min(max_ch, Inches(need + 1.02))
    used = ch * rows + gap * (rows - 1)
    # nudge the block down a little, but never let it float mid-slide
    top += min(Inches(0.34), max(0, (avail - used) // 2))

    for i, item in enumerate(items):
        r, c = divmod(i, cols)
        x = MARGIN + c * (cw + gap)
        y = top + r * (ch + gap)
        rect(s, x, y, cw, ch, fill=WHITE, outline=LINE, radius=0.05)
        rect(s, x, y, cw, Inches(0.045), fill=accent, radius=0)
        ty = y + Inches(0.26)
        if numbered:
            circle(s, x + Inches(0.34), ty + Inches(0.11), Inches(0.34), fill=wash,
                   label=str(i + 1), size=12, color=accent)
            tx = x + Inches(0.62)
        else:
            tx = x + Inches(0.28)
        _, tf = textbox(s, tx, ty, cw - (tx - x) - Inches(0.24), Inches(0.4))
        para(tf, item[0], size=fit_size(item[0], (cw - (tx - x) - Inches(0.24)) / 914400, 0.38,
                                        max_pt=13, min_pt=10.6, line=1.14,
                                        font=HEAD_FONT, bold=True),
             color=INK, font=HEAD_FONT, bold=True, first=True, space_after=0, line=1.14)
        bw, bh = cw - Inches(0.52), ch - Inches(0.86)
        b_pt = fit_size(item[1], bw / 914400, bh / 914400, max_pt=body_size, min_pt=7.8, line=1.3)
        _, tf = textbox(s, x + Inches(0.28), ty + Inches(0.52), bw, bh)
        para(tf, item[1], size=b_pt, color=MUTED, first=True, space_after=0, line=1.3)
    return s


def matrix(prs, eb, title, sub, headers, rows, note=None, widths=None, accent=VIOLET, path=None):
    s = page(prs, eb, title, sub, note, path)
    top = body_top({"path": path, "sub": sub}) + Inches(0.04)
    widths = widths or [1.0 / len(headers)] * len(headers)
    total = CONTENT_W
    xs, acc = [], MARGIN
    for wgt in widths:
        xs.append(acc)
        acc += int(total * wgt)

    hh = Inches(0.44)
    rect(s, MARGIN, top, total, hh, fill=WASH, radius=0.04)
    for i, head in enumerate(headers):
        _, tf = textbox(s, xs[i] + Inches(0.18), top + Inches(0.11),
                        int(total * widths[i]) - Inches(0.3), Inches(0.3))
        para(tf, head.upper(), size=9.5, color=accent, font=HEAD_FONT, bold=True, first=True,
             space_after=0)

    y = top + hh + Inches(0.1)
    rh = (H - y - Inches(0.74) - Inches(0.08) * (len(rows) - 1)) // len(rows)
    for r, row in enumerate(rows):
        ry = y + r * (rh + Inches(0.08))
        rect(s, MARGIN, ry, total, rh, fill=CANVAS, outline=LINE, radius=0.05)
        for i, cell in enumerate(row):
            _, tf = textbox(s, xs[i] + Inches(0.18), ry + Inches(0.13),
                            int(total * widths[i]) - Inches(0.32), rh - Inches(0.24))
            bold = i == 0
            cw_in = (int(total * widths[i]) - Inches(0.32)) / 914400
            pt = fit_size(cell, cw_in, (rh - Inches(0.26)) / 914400,
                          max_pt=11.2 if bold else 10.5, min_pt=7.8, line=1.3,
                          font=HEAD_FONT if bold else BODY_FONT, bold=bold)
            para(tf, cell, size=pt,
                 color=INK if bold else MUTED,
                 font=HEAD_FONT if bold else BODY_FONT, bold=bold, first=True,
                 space_after=0, line=1.3)
    return s


def split(prs, eb, title, sub, left, right, note=None, path=None):
    """Two contrasting columns — typically 'works today' against 'still to agree'."""
    s = page(prs, eb, title, sub, note, path)
    top = body_top({"path": path, "sub": sub}) + Inches(0.04)
    cw = (CONTENT_W - Inches(0.28)) // 2
    for idx, (label, colour, wash, items) in enumerate([left, right]):
        x = MARGIN + idx * (cw + Inches(0.28))
        rect(s, x, top, cw, H - top - Inches(0.74), fill=WHITE, outline=LINE, radius=0.04)
        rect(s, x, top, cw, Inches(0.05), fill=colour, radius=0)
        badge(s, x + Inches(0.26), top + Inches(0.26), label, fill=wash, color=colour, size=10)
        y = top + Inches(0.76)
        avail = (H - Inches(0.86)) - y
        slot = avail // len(items)
        iw = (cw - Inches(0.52)) / 914400
        for t, b in items:
            _, tf = textbox(s, x + Inches(0.26), y, cw - Inches(0.52), Inches(0.3))
            para(tf, t, size=12.4, color=INK, font=HEAD_FONT, bold=True, first=True,
                 space_after=0, line=1.15)
            bh = slot - Inches(0.36)
            pt = fit_size(b, iw, bh / 914400, max_pt=10.6, min_pt=8.0, line=1.3)
            _, tf = textbox(s, x + Inches(0.26), y + Inches(0.28), cw - Inches(0.52), bh)
            para(tf, b, size=pt, color=MUTED, first=True, space_after=0, line=1.3)
            y += slot
    return s


def stats(prs, eb, title, sub, tiles, note=None, body=None, path=None):
    """A row of headline numbers, optionally over a panel of short prose."""
    s = page(prs, eb, title, sub, note, path)
    top = body_top({"path": path, "sub": sub}) + Inches(0.12)
    n = len(tiles)
    gap = Inches(0.18)
    cw = (CONTENT_W - gap * (n - 1)) // n
    tile_h = Inches(1.92)
    inner = cw - Inches(0.48)
    for i, (value, label, detail) in enumerate(tiles):
        x = MARGIN + i * (cw + gap)
        rect(s, x, top, cw, tile_h, fill=WASH, radius=0.05)
        _, tf = textbox(s, x + Inches(0.24), top + Inches(0.2), inner, Inches(0.62))
        para(tf, value, size=34, color=VIOLET, font=HEAD_FONT, bold=True, first=True,
             space_after=0, line=1.0)
        lbl_pt = fit_size(label, inner / 914400, 0.44, max_pt=11.5, min_pt=9.4, line=1.14,
                          font=HEAD_FONT, bold=True)
        _, tf = textbox(s, x + Inches(0.24), top + Inches(0.86), inner, Inches(0.46))
        para(tf, label, size=lbl_pt, color=INK, font=HEAD_FONT, bold=True, first=True,
             space_after=0, line=1.14)
        _, tf = textbox(s, x + Inches(0.24), top + Inches(1.34), inner, Inches(0.5))
        para(tf, detail, size=fit_size(detail, inner / 914400, 0.48, max_pt=9.6, min_pt=7.8,
                                       line=1.24),
             color=MUTED, first=True, space_after=0, line=1.24)

    if body:
        by = top + tile_h + Inches(0.26)
        avail = H - by - Inches(0.74)
        panel = rect(s, MARGIN, by, CONTENT_W, avail, fill=CANVAS, outline=LINE, radius=0.03)
        tw = (CONTENT_W - Inches(0.6)) / 914400
        heights, sizes = [], []
        for t, b in body:
            pt = fit_size(b, tw, 1.2, max_pt=11.0, min_pt=9.4, line=1.34)
            sizes.append(pt)
            heights.append(0.3 + text_lines(b, tw, pt) * pt * 1.34 / 72)
        total = sum(heights) + 0.16 * (len(body) - 1)
        y = by + Inches(max(0.24, (avail / 914400 - total) / 2))
        for i, (t, b) in enumerate(body):
            _, tf = textbox(s, MARGIN + Inches(0.3), y, CONTENT_W - Inches(0.6), Inches(0.28))
            para(tf, t, size=12.6, color=INK, font=HEAD_FONT, bold=True, first=True,
                 space_after=0, line=1.2)
            _, tf = textbox(s, MARGIN + Inches(0.3), y + Inches(0.28),
                            CONTENT_W - Inches(0.6), Inches(heights[i] - 0.28))
            para(tf, b, size=sizes[i], color=MUTED, first=True, space_after=0, line=1.34)
            y += Inches(heights[i] + 0.16)
    return s


def gif_slide(prs, eb, title, sub, gif_path, steps, note=None, path=None):
    s = page(prs, eb, title, sub, note, path)
    top = body_top({"path": path, "sub": sub}) + Inches(0.04)
    pic_w = Inches(8.05)
    iw, ih = Image.open(gif_path).size
    scale = min(pic_w / iw, (H - top - Inches(0.78)) / ih)
    s.shapes.add_picture(gif_path, MARGIN, top, int(iw * scale), int(ih * scale))
    cx = MARGIN + pic_w + Inches(0.32)
    cw = W - cx - MARGIN
    y = top
    for i, (t, b) in enumerate(steps):
        circle(s, cx + Inches(0.17), y + Inches(0.17), Inches(0.34), fill=VIOLET,
               label=str(i + 1), size=12)
        _, tf = textbox(s, cx + Inches(0.44), y + Inches(0.02), cw - Inches(0.5), Inches(0.3))
        para(tf, t, size=12.2, color=INK, font=HEAD_FONT, bold=True, first=True,
             space_after=0, line=1.15)
        _, tf = textbox(s, cx + Inches(0.44), y + Inches(0.3), cw - Inches(0.5), Inches(0.62))
        para(tf, b, size=fit_size(b, (cw - Inches(0.5)) / 914400, 0.6,
                                  max_pt=10.2, min_pt=8.2, line=1.28),
             color=MUTED, first=True, space_after=0, line=1.28)
        y += Inches(0.92)
    tag = rect(s, cx, H - Inches(1.42), cw, Inches(0.56), fill=WASH, radius=0.08)
    tf = tag.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.18)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "ANIMATED CAPTURE  ·  PLAYS IN SLIDE SHOW"
    r.font.size = Pt(10.5)
    r.font.bold = True
    r.font.name = HEAD_FONT
    r.font.color.rgb = VIOLET
    return s


def closing(prs, title, items, footer):
    s = blank(prs)
    page_bg(s, INK)
    rect(s, 0, 0, Inches(0.16), H, fill=VIOLET_2, radius=0)
    _, tf = textbox(s, MARGIN + Inches(0.3), Inches(0.92), Inches(10.4), Inches(0.9))
    para(tf, title, size=34, color=WHITE, font=HEAD_FONT, bold=True, first=True,
         space_after=0, line=1.06)
    y = Inches(2.06)
    for i, (t, b) in enumerate(items):
        circle(s, MARGIN + Inches(0.48), y + Inches(0.24), Inches(0.46), fill=VIOLET_2,
               label=str(i + 1), size=15)
        _, tf = textbox(s, MARGIN + Inches(1.02), y, Inches(10.6), Inches(0.32))
        para(tf, t, size=15, color=WHITE, font=HEAD_FONT, bold=True, first=True,
             space_after=0, line=1.15)
        _, tf = textbox(s, MARGIN + Inches(1.02), y + Inches(0.32), Inches(10.4), Inches(0.6))
        para(tf, b, size=fit_size(b, 10.4, 0.58, max_pt=11.6, min_pt=9.6, line=1.34),
             color=RGB_SOFT(), first=True, space_after=0, line=1.34)
        y += Inches(1.12)
    _, tf = textbox(s, MARGIN + Inches(0.3), H - Inches(0.72), Inches(11.4), Inches(0.4))
    para(tf, footer, size=10, color=VIOLET_2, first=True, space_after=0, italic=True)
    return s
