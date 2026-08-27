"""Deck design system: palette, type scale and the layout primitives every
slide is built from.

The palette is taken from the prototype's own CSS tokens so screenshots sit on
the slide without a colour clash.
"""
import json
import os

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ------------------------------------------------------------------ palette
#
# Defaults are sampled from the prototype's own CSS tokens, so a screenshot sits
# on the slide without a colour clash. Drop a brand.json next to the build
# scripts to override any of them:
#
#   { "INK": "241243", "VIOLET": "54248F", "HEAD_FONT": "Segoe UI Semibold" }
#
_DEFAULTS = {
    "INK": "241243",        # deep heading colour
    "VIOLET": "54248F",     # primary brand
    "VIOLET_2": "7C3FC2",   # lighter brand, used on dark pages
    "WASH": "F2EBFA",       # tinted panel
    "CANVAS": "F7F8FB",     # card fill
    "LINE": "E8E7ED",       # hairline border
    "MUTED": "6B6675",      # body text
    "WHITE": "FFFFFF",
    "MARK": "E4572E",       # callout markers, on-screen and on-slide
    "MARK_WASH": "FDEEE8",
    "LIVE": "0F766E",       # "shown on a live screen"
    "LIVE_WASH": "E6F4F2",
    "DECIDE": "B45309",     # "for the meeting"
    "DECIDE_WASH": "FDF3E6",
}
_FONTS = {
    "HEAD_FONT": "Segoe UI Semibold",
    "BODY_FONT": "Segoe UI",
    "MONO_FONT": "Consolas",
}

_brand = {}
for _candidate in ("brand.json", os.path.join("..", "brand.json")):
    if os.path.exists(_candidate):
        with open(_candidate, encoding="utf-8") as _fh:
            _brand = json.load(_fh)
        break

for _key, _hex in _DEFAULTS.items():
    _value = str(_brand.get(_key, _hex)).lstrip("#")
    globals()[_key] = RGBColor(int(_value[0:2], 16), int(_value[2:4], 16), int(_value[4:6], 16))
for _key, _name in _FONTS.items():
    globals()[_key] = _brand.get(_key, _name)

INK = globals()["INK"]; VIOLET = globals()["VIOLET"]; VIOLET_2 = globals()["VIOLET_2"]
WASH = globals()["WASH"]; CANVAS = globals()["CANVAS"]; LINE = globals()["LINE"]
MUTED = globals()["MUTED"]; WHITE = globals()["WHITE"]; MARK = globals()["MARK"]
MARK_WASH = globals()["MARK_WASH"]; LIVE = globals()["LIVE"]; LIVE_WASH = globals()["LIVE_WASH"]
DECIDE = globals()["DECIDE"]; DECIDE_WASH = globals()["DECIDE_WASH"]
HEAD_FONT = globals()["HEAD_FONT"]; BODY_FONT = globals()["BODY_FONT"]
MONO_FONT = globals()["MONO_FONT"]

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)
CONTENT_W = W - MARGIN * 2


# ------------------------------------------------------------- text fitting
# Files used only for *measuring* text so the builder can pick a size that fits.
# PowerPoint renders with the family names above; these just need to be close.
_FIT_FONTS = {
    ("Segoe UI", False): r"C:\Windows\Fonts\segoeui.ttf",
    ("Segoe UI", True): r"C:\Windows\Fonts\segoeuib.ttf",
    ("Segoe UI Semibold", False): r"C:\Windows\Fonts\seguisb.ttf",
    ("Segoe UI Semibold", True): r"C:\Windows\Fonts\seguisb.ttf",
    ("Calibri", False): r"C:\Windows\Fonts\calibri.ttf",
    ("Calibri", True): r"C:\Windows\Fonts\calibrib.ttf",
    ("Arial", False): r"C:\Windows\Fontsrial.ttf",
    ("Arial", True): r"C:\Windows\Fontsrialbd.ttf",
}
_FIT_FONTS = {k: v for k, v in _FIT_FONTS.items() if os.path.exists(v)}
if not _FIT_FONTS:
    raise RuntimeError("No measuring font found; add one to _FIT_FONTS in theme.py")
_FALLBACK = next(iter(_FIT_FONTS.values()))
_fit_cache = {}
_measure = None


def _fit_font(name, pts, bold):
    from PIL import ImageFont
    key = (name, bold, round(pts * 4))
    if key not in _fit_cache:
        path = _FIT_FONTS.get((name, bold)) or _FALLBACK
        _fit_cache[key] = ImageFont.truetype(path, max(4, round(pts * 4)))
    return _fit_cache[key]


def text_lines(text, width_in, pts, font=BODY_FONT, bold=False):
    """How many wrapped lines `text` needs in a box `width_in` inches wide.

    Measured at 4x with real font metrics, which tracks PowerPoint closely
    enough to choose a size that will not overflow.
    """
    global _measure
    if _measure is None:
        from PIL import Image, ImageDraw
        _measure = ImageDraw.Draw(Image.new("RGB", (8, 8)))
    fnt = _fit_font(font, pts, bold)
    limit = width_in * 72 * 4          # inches -> quarter-points, matching the 4x font
    lines, line = 1, ""
    for word in text.split():
        probe = (line + " " + word).strip()
        if _measure.textlength(probe, font=fnt) <= limit or not line:
            line = probe
        else:
            lines += 1
            line = word
    return lines


def text_width(text, pts, font=BODY_FONT, bold=False):
    """Rendered width of a single line, in inches."""
    global _measure
    if _measure is None:
        from PIL import Image, ImageDraw
        _measure = ImageDraw.Draw(Image.new("RGB", (8, 8)))
    return _measure.textlength(text, font=_fit_font(font, pts, bold)) / (72 * 4)


def fit_size(text, width_in, height_in, max_pt=11.0, min_pt=8.0, line=1.28,
             font=BODY_FONT, bold=False, step=0.25):
    """Largest point size at which `text` still fits the given box."""
    size = max_pt
    while size > min_pt:
        n = text_lines(text, width_in, size, font=font, bold=bold)
        if n * size * line / 72 <= height_in:
            return size
        size -= step
    return min_pt


def emu_in(value):
    return value / 914400


# ------------------------------------------------------------------- helpers
def textbox(slide, x, y, w, h, wrap=True, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


def para(tf, text, size=14, color=MUTED, font=BODY_FONT, bold=False, first=False,
         space_before=0, space_after=6, line=1.32, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def rect(slide, x, y, w, h, fill=WHITE, outline=None, radius=0.055, shadow=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if radius:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if outline is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = outline
        shape.line.width = Pt(1)
    if not shadow:
        shape.shadow.inherit = False
    shape.text_frame.word_wrap = True
    return shape


def circle(slide, cx, cy, d, fill=MARK, label="", size=13, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - d // 2, cy - d // 2, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(label)
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = HEAD_FONT
    r.font.color.rgb = color
    return shape


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def page_bg(slide, color=WHITE):
    shape = rect(slide, 0, 0, W, H, fill=color, radius=0)
    shape.shadow.inherit = False
    return shape


def eyebrow(slide, text, y=Inches(0.46), color=VIOLET_2):
    _, tf = textbox(slide, MARGIN, y, CONTENT_W, Inches(0.26))
    para(tf, text.upper(), size=10.5, color=color, font=HEAD_FONT, bold=True, first=True,
         space_after=0)


def heading(slide, text, y=Inches(0.76), size=29, color=INK, width=None):
    w = width or CONTENT_W
    size = fit_size(text, w / 914400, 0.68, max_pt=size, min_pt=20, line=1.05,
                    font=HEAD_FONT, bold=True)
    _, tf = textbox(slide, MARGIN, y, w, Inches(0.72))
    para(tf, text, size=size, color=color, font=HEAD_FONT, bold=True, first=True,
         space_after=0, line=1.05)


def subhead(slide, text, y=Inches(1.42), size=13.5, color=MUTED, width=None):
    _, tf = textbox(slide, MARGIN, y, width or Inches(10.6), Inches(0.5))
    para(tf, text, size=size, color=color, first=True, space_after=0, line=1.28)


def path_chip(slide, text, y=Inches(1.44)):
    """The navigation path to a screen, rendered as one wash-filled pill.

    Every screen slide carries one so a reader can reach the same page in ATLAS
    without being told separately.
    """
    size = fit_size(text, 11.4, 0.24, max_pt=10.4, min_pt=8.2, line=1.0,
                    font=HEAD_FONT, bold=True)
    w = min(CONTENT_W, Inches(text_width(text, size, HEAD_FONT, True) + 0.52))
    shape = rect(slide, MARGIN, y, w, Inches(0.34), fill=WASH, radius=0.5)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.22)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = HEAD_FONT
    r.font.color.rgb = VIOLET
    return shape


def footnote(slide, text, y=None):
    y = y or (H - Inches(0.52))
    _, tf = textbox(slide, MARGIN, y, CONTENT_W, Inches(0.26))
    para(tf, text, size=9, color=MUTED, first=True, space_after=0, italic=True)


def badge(slide, x, y, text, fill=LIVE_WASH, color=LIVE, w=None, h=Inches(0.28), size=9.5):
    w = w or Inches(0.14 * len(text) + 0.24)
    shape = rect(slide, x, y, w, h, fill=fill, radius=0.5)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = HEAD_FONT
    r.font.color.rgb = color
    return shape


def picture_fit(slide, path, x, y, max_w, max_h, align="center", valign="center"):
    """Place an image scaled to fit a box, keeping its aspect ratio."""
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    if align == "center":
        px = x + (max_w - w) // 2
    elif align == "left":
        px = x
    else:
        px = x + max_w - w
    py = y if valign == "top" else y + (max_h - h) // 2
    return slide.shapes.add_picture(path, int(px), int(py), int(w), int(h))


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text
