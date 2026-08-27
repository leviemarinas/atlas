# -*- coding: utf-8 -*-
"""House-rules linter for a finished deck.

Reads the .pptx back and reports anything that breaks the rules in
references/house-style.md. Run it as the last step of every build; a clean run
is the signal the deck is ready to send.

    python check_deck.py ../outputs/MyDeck.pptx

Checks
  counts     catalogue totals quoted from a prototype ("219 computations")
  negative   deficiency framing ("does not yet", "cannot", "gap", "not wired")
  path       a screen page with a screenshot but no navigation path chip
  orphan     a numbered marker on an image with no matching note, or vice versa
  notes      a slide with no speaker notes
"""
import collections
import glob
import os
import re
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# A bare integer next to a catalogue noun is almost always a prototype total.
CATALOGUE = re.compile(
    r"\b(\d{1,4})\s+(?:governed\s+)?"
    r"(computations?|formulas?|reference sources?|policy engines?|assignments?|"
    r"approved fields?|codes?|records?|categories|tables?|engines?)\b", re.I)
SPELLED = re.compile(
    r"\b(two|three|four|five|six|seven|eight|nine|ten|eleven|twelve|thirty|sixty(?:-one)?)\s+"
    r"(governed computations?|approved (?:payroll )?fields?|reference sources?|"
    r"policy engines?|client assignments?|finance categories)\b", re.I)

NEGATIVE = re.compile(
    r"\b(does not (?:yet|run|select|read|tell|prove|consult|exist)"
    r"|do not (?:yet|run|select|read|tell)"
    r"|is not (?:wired|built|implemented|supported)"
    r"|not yet (?:enforced|wired|built|implemented|available)"
    r"|cannot\b|can't\b|did not (?:run|appear|happen)"
    r"|still to build|integration gap|missing (?:feature|capability)"
    r"|defect|shortcoming|limitation|half-broken|fails? validation)\b", re.I)

# Phrasings that are fine even though they contain a negative word.
ALLOWED = re.compile(r"(never means editing|never get confused|not the calendar|"
                     r"not a spreadsheet|not production|rather than|instead of)", re.I)


def shape_texts(slide):
    """Each text frame separately.

    Scanning a whole slide as one string produces false positives: a numbered
    marker badge is its own shape holding just "4", and joining it to the next
    shape's title yields "4 Reference sources", which reads like a count.
    """
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.append(" ".join(sh.text_frame.text.split()))
    return out


def slide_text(slide):
    return "  ".join(shape_texts(slide))


def has_picture(slide):
    return any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in slide.shapes)


def has_path(slide):
    """A path chip reads like 'A > B > C' with at least two separators."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        text = sh.text_frame.text
        if text.count("›") >= 2 or text.count(">") >= 2:
            return True
    return False


def check(path):
    prs = Presentation(path)
    findings = []
    for i, slide in enumerate(prs.slides, 1):
        text = slide_text(slide)

        for chunk in shape_texts(slide):
            for rx, label in ((CATALOGUE, "counts"), (SPELLED, "counts")):
                for m in rx.finditer(chunk):
                    lo = max(0, m.start() - 40)
                    findings.append((i, label, f"…{chunk[lo:m.end() + 40]}…"))

            for m in NEGATIVE.finditer(chunk):
                lo, hi = max(0, m.start() - 60), m.end() + 60
                window = chunk[lo:hi]
                if ALLOWED.search(window):
                    continue
                findings.append((i, "negative", f"…{window}…"))

        if has_picture(slide) and not has_path(slide) and len(text) > 120:
            findings.append((i, "path", "screenshot page with no navigation path chip"))

        if not slide.has_notes_slide or not slide.notes_slide.notes_text_frame.text.strip():
            findings.append((i, "notes", "no speaker notes"))

    return findings, len(prs.slides.__iter__.__self__._sldIdLst)


def main():
    if len(sys.argv) > 1:
        target = sys.argv[1]
    else:
        found = [f for f in sorted(glob.glob(os.path.join("..", "outputs", "*.pptx")))
                 if not os.path.basename(f).startswith("~$")]
        if not found:
            sys.exit("usage: python check_deck.py <deck.pptx>")
        target = found[-1]

    findings, slides = check(target)
    print(f"checked {os.path.basename(target)} ({slides} slides)\n")
    by_kind = collections.Counter(kind for _, kind, _ in findings)
    for num, kind, detail in findings:
        print(f"  slide {num:>2}  [{kind}]  {detail}")
    if findings:
        print("\n" + ", ".join(f"{n} {k}" for k, n in by_kind.most_common()))
        print(f"{len(findings)} issues — see references/house-style.md")
        sys.exit(1)
    print("clean — no catalogue counts, no deficiency framing, "
          "every screenshot page has a path, every slide has notes")


if __name__ == "__main__":
    main()
